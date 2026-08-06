#!/usr/bin/env python3
"""
generate_presentacion_comite.py
Genera la presentación PPTX para la reunión de expertos: socialización de
Antioquia Natural (explicación no técnica) + criterios de evaluación y
análisis de selección de especies (158 especies propuestas).

Fuente de los datos: generate_evaluacion_especies.js (7 criterios, umbrales
de categoría) y generate_analisis_especies_doc.js (propuesta de 158
especies, distribución por grupo, especies insignia, cobertura geográfica,
recomendaciones al comité); no se inventó ningún dato.

Uso: python3 backend/src/scripts/generate_presentacion_comite.py
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.join(os.path.dirname(__file__), '../../../')
OUT_DIR = os.path.join(ROOT, 'Documentos gobernacion')
OUT_FILE = os.path.join(OUT_DIR, 'Presentacion_Comite_Cientifico_Antioquia_Natural.pptx')
LOGO = os.path.join(ROOT, 'biodiversidad/img/logo/logo_gobernacion.png')

# ── Colores institucionales ──────────────────────────────────────────────
GREEN = RGBColor(0x01, 0x8D, 0x38)
DARK_GREEN = RGBColor(0x0B, 0x56, 0x40)
GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

PRIORITARIA = RGBColor(0xA5, 0xD6, 0xA7)
RECOMENDADA = RGBColor(0xFF, 0xF5, 0x9D)
CONDICIONAL = RGBColor(0xFF, 0xCC, 0x80)
ESPERA = RGBColor(0xEF, 0x9A, 0x9A)

FONT = 'Arial'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False, font=FONT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK_TEXT,
                 space_after=10, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(space_after)
        p.line_spacing = 1.05
        bullet_char = '●  ' if level == 0 else '–  '
        segs = text.split('**')
        first = True
        for j, seg in enumerate(segs):
            run = p.add_run()
            run.text = (bullet_char if first else '') + seg
            first = False
            run.font.size = Pt(size - level * 1)
            run.font.bold = (j % 2 == 1)
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_header(slide, kicker, title):
    """Barra superior verde con kicker + título de la diapositiva."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = kicker.upper()
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = LIGHT_GREEN
    r.font.name = FONT
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = title
    r2.font.size = Pt(26)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    r2.font.name = FONT


def add_footer(slide, n):
    add_textbox(slide, Inches(0.5), Inches(7.15), Inches(6), Inches(0.3),
                'Antioquia Natural · Reunión de Expertos', size=10, color=GRAY_TEXT)
    add_textbox(slide, SLIDE_W - Inches(1.2), Inches(7.15), Inches(0.8), Inches(0.3),
                str(n), size=10, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)


def content_slide(kicker, title, n):
    s = add_slide()
    set_bg(s)
    add_header(s, kicker, title)
    add_footer(s, n)
    return s


def add_table(slide, left, top, width, height, headers, rows, col_widths=None,
              font_size=12, header_size=13, row_colors=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gtable.table
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(width * (w / total))
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(header_size)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = FONT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if row_colors and r in row_colors:
                cell.fill.fore_color.rgb = row_colors[r]
            else:
                cell.fill.fore_color.rgb = LIGHT_GREEN if r % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.color.rgb = DARK_TEXT
                run.font.name = FONT
    return gtable


def section_divider(n_text, title, subtitle):
    s = add_slide()
    set_bg(s, DARK_GREEN)
    add_textbox(s, Inches(1), Inches(2.2), Inches(3), Inches(1), n_text,
                size=20, bold=True, color=RGBColor(0x9C, 0xCC, 0xB0), font=FONT)
    add_textbox(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.9), title,
                size=38, bold=True, color=WHITE, font=FONT, line_spacing=1.05)
    add_textbox(s, Inches(1), Inches(4.9), Inches(11.3), Inches(1), subtitle,
                size=18, color=RGBColor(0xC8, 0xE6, 0xC9), font=FONT)
    return s


def card(slide, left, top, width, height, emoji_title, body, fill=LIGHT_GREEN):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = GREEN
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = emoji_title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = DARK_GREEN
    r.font.name = FONT
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(13)
    r2.font.color.rgb = DARK_TEXT
    r2.font.name = FONT


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1: PORTADA
# ══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, DARK_GREEN)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.4), Inches(0.7), height=Inches(1.3))
add_textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
            'Antioquia Natural', size=44, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font=FONT)
add_textbox(s, Inches(1), Inches(3.6), Inches(11.3), Inches(1),
            'Reunión de Expertos: Criterios de Evaluación y Análisis de Selección de Especies',
            size=22, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.CENTER, font=FONT)
add_textbox(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
            'Gobernación de Antioquia · Secretaría de Ambiente', size=14,
            color=RGBColor(0x9C, 0xCC, 0xB0), align=PP_ALIGN.CENTER, font=FONT)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ══════════════════════════════════════════════════════════════════════════
s = content_slide('Agenda', 'De qué vamos a hablar hoy (20 minutos)', 2)
card(s, Inches(0.7), Inches(1.8), Inches(5.7), Inches(4.6),
     'Parte 1\n¿Qué es Antioquia Natural?',
     '\nUna introducción breve a la aplicación: para qué sirve, cómo se usa y qué contiene cada uno de sus módulos.\n\n(≈ 8 minutos)')
card(s, Inches(6.9), Inches(1.8), Inches(5.7), Inches(4.6),
     'Parte 2\nCriterios y propuesta de especies',
     '\nCómo evaluamos qué especies incluir, y la propuesta de 158 especies lista para discusión del comité.\n\n(≈ 12 minutos)')

# ══════════════════════════════════════════════════════════════════════════
# PARTE 1
# ══════════════════════════════════════════════════════════════════════════
section_divider('PARTE 1 DE 2', '¿Qué es Antioquia Natural?',
                 'Una introducción no técnica a la aplicación')

# SLIDE 4: Qué es / cómo se usa
s = content_slide('Parte 1', '¿Qué es y cómo se usa?', 4)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5),
            [
                ('Antioquia Natural es una aplicación web para **consultar la biodiversidad, el agua y los programas comunitarios del departamento**, pensada para el público general.', 0),
                ('Se accede escaneando un **código QR** con el celular: no hay que instalar nada.', 0),
                ('Funciona en **español e inglés**, pensada para uso en campo, colegios, senderos ecológicos y puntos de información turística.', 0),
                ('Diseñada primero para **celular** (mobile-first): botones grandes, fotos, poco texto por pantalla.', 0),
            ], size=19, space_after=22)

# SLIDE 5: Los 3 módulos overview
s = content_slide('Parte 1', 'Tres módulos principales', 5)
card(s, Inches(0.6), Inches(1.7), Inches(3.9), Inches(4.5), '🌿 Biodiversidad',
     '\nCatálogo de flora y fauna del departamento, organizado por subregión y grupo (aves, mariposas, árboles, etc.)')
card(s, Inches(4.7), Inches(1.7), Inches(3.9), Inches(4.5), '💧 Agua',
     '\nMapa de la red hídrica: cuencas principales y fuentes de agua que abastecen acueductos municipales.')
card(s, Inches(8.8), Inches(1.7), Inches(3.9), Inches(4.5), '🤝 Comunidad',
     '\nProgramas ciudadanos: Jóvenes pa\' Lante y Guarda Cuencas: jóvenes y comunidades documentando su territorio con fotos.')

# SLIDE 6: Módulo Biodiversidad en detalle
s = content_slide('Parte 1', 'Módulo Biodiversidad: qué encuentra el usuario', 6)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5),
            [
                ('El usuario navega por **9 grupos**: aves, anfibios y reptiles, mariposas, polillas, mamíferos, animales domésticos, peces de agua dulce, orquídeas y árboles nativos.', 0),
                ('Puede buscar por **subregión** (las 9 subregiones de Antioquia) o directamente por nombre de especie.', 0),
                ('Cada especie tiene una **ficha**: fotos, nombre común y científico, categoría de conservación (IUCN), distribución y una descripción bilingüe.', 0),
                ('Es exactamente este catálogo, qué especies incluir y con qué criterio, el tema central de la reunión de hoy.', 0),
            ], size=19, space_after=22)

# SLIDE 7: Comunidad
s = content_slide('Parte 1', 'Módulo Comunidad: ciencia ciudadana', 7)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5),
            [
                ('**Jóvenes pa\' Lante**: jóvenes de 90 municipios suben fotos de biodiversidad observada en su territorio cada mes.', 0),
                ('**Guarda Cuencas**: comunidades documentan el estado de las cuencas hídricas de su región.', 0),
                ('**Especie del Mes**: una especie destacada mensualmente, con galería de avistamientos ciudadanos.', 0),
                ('Este módulo conecta directamente con el catálogo de especies: entre más completo y confiable sea el catálogo, más valor tiene lo que la comunidad reporta.', 0),
            ], size=19, space_after=22)

# SLIDE 8: Transición
s = content_slide('Parte 1 → Parte 2', 'El catálogo de especies es el corazón de la app', 8)
add_textbox(s, Inches(0.9), Inches(2.2), Inches(11), Inches(3.5),
            'Todo lo que acabamos de ver (el módulo de Biodiversidad, la Especie del Mes, lo que suben los jóvenes de los programas) depende de una sola decisión: qué especies incluimos en el catálogo, y con qué criterio.\n\nDe eso trata la segunda parte de esta reunión.',
            size=24, color=DARK_TEXT, line_spacing=1.3)

# ══════════════════════════════════════════════════════════════════════════
# PARTE 2
# ══════════════════════════════════════════════════════════════════════════
section_divider('PARTE 2 DE 2', 'Criterios de Evaluación y Selección de Especies',
                 'Cómo decidimos qué especies incluir en el catálogo')

# SLIDE 10: El reto
s = content_slide('Parte 2', 'El reto: ¿cómo elegir qué especies incluir?', 10)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5),
            [
                ('Antioquia tiene una biodiversidad enorme: no es posible ni deseable incluir todo desde el primer lanzamiento.', 0),
                ('Necesitábamos un método **objetivo y verificable**, no una selección a criterio personal, que el comité pudiera revisar y ajustar.', 0),
                ('La solución: un sistema de **puntaje** con 7 criterios técnicos, aplicado especie por especie.', 0),
            ], size=20, space_after=26)

# SLIDE 11: Metodología overview
s = content_slide('Parte 2', 'La metodología: 7 criterios, puntaje de 0 a 100', 11)
add_table(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(4.6),
          ['Criterio', 'Qué mide', 'Puntaje máximo'],
          [
              ['P1 · Estado de conservación (IUCN)', 'Qué tan amenazada está la especie', '20'],
              ['P2 · Endemismo', 'Si es exclusiva de Antioquia, los Andes o Colombia', '20'],
              ['P3 · Rol ecológico', 'Qué tan insustituible es en su ecosistema', '15'],
              ['P4 · Representación geográfica', 'Si cubre subregiones poco representadas', '15'],
              ['P5 · Atractivo (reconocibilidad + historia + valor cultural)', 'Qué tan fácil es que conecte con el público', '15'],
              ['P6 · Disponibilidad de foto', 'Qué tan fácil es conseguir una foto de calidad', '10'],
              ['P7 · Diversidad taxonómica', 'Si evita repetir familias ya cubiertas', '5'],
          ], col_widths=[4.5, 5.5, 1.9], font_size=13, header_size=14)

# SLIDE 12: De puntaje a categoría
s = content_slide('Parte 2', 'De puntaje a decisión: 4 categorías', 12)
add_table(s, Inches(1.5), Inches(1.7), Inches(9.9), Inches(3.6),
          ['Puntaje', 'Categoría', 'Decisión recomendada'],
          [
              ['≥ 80', 'Prioritaria', 'Debe estar en la versión 1.0 sin excepción'],
              ['60 – 79', 'Recomendada', 'Incluir dentro del cupo disponible del grupo'],
              ['40 – 59', 'Condicional', 'Incluir si hay cupo y cumple cuota mínima'],
              ['< 40', 'Lista de espera', 'Reservar para Fase 2 o reemplazos'],
          ], col_widths=[1.8, 2.6, 5.5], font_size=15, header_size=15,
          row_colors={0: PRIORITARIA, 1: RECOMENDADA, 2: CONDICIONAL, 3: ESPERA})
add_textbox(s, Inches(1.5), Inches(5.6), Inches(9.9), Inches(1),
            'Fuentes de información: Lista Roja de la IUCN, Libros Rojos de Colombia (Instituto Humboldt), SiB Colombia y literatura científica sobre biodiversidad antioqueña.',
            size=13, italic=True, color=GRAY_TEXT)

# SLIDE 13: La propuesta: números generales
s = content_slide('Parte 2', 'La propuesta: 158 especies para discusión del comité', 13)
add_textbox(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.6),
            'Hoy la app tiene 117 especies documentadas en vivo. La propuesta que sometemos a discusión del comité expande el catálogo a 158.',
            size=16, color=GRAY_TEXT)
add_table(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.3),
          ['Grupo', 'Propuestas', 'Prioritarias (≥80)', 'Con foto lista (A/B)'],
          [
              ['Aves', '30', '8', '26'],
              ['Anfibios y Reptiles', '25', '7', '18'],
              ['Mariposas', '20', '4', '17'],
              ['Orquídeas', '20', '6', '15'],
              ['Árboles Nativos', '18', '5', '14'],
              ['Peces de Agua Dulce', '15', '6', '9'],
              ['Mamíferos', '13', '5', '10'],
              ['Polillas', '10', '2', '8'],
              ['Animales Domésticos', '7', '3', '6'],
              ['Total', '158', '46', '123'],
          ], col_widths=[4, 2, 2.5, 3], font_size=13, header_size=14,
          row_colors={9: PRIORITARIA})

# SLIDE 14: Especies insignia
s = content_slide('Parte 2', 'Especies insignia: los casos más urgentes', 14)
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.9),
          ['Especie', 'Grupo', 'Puntaje', 'Por qué es urgente'],
          [
              ['Prochilodus magdalenae', 'Peces', '86', 'CR + especie clave de la pesca artesanal'],
              ['Ateles hybridus', 'Mamíferos', '85', 'CR + uno de los 25 primates más amenazados del mundo'],
              ['Crax alberti', 'Aves', '82', 'CR + endémica de Colombia'],
              ['Lipaugus weberi', 'Aves', '78', 'EN + endémica exclusiva de Antioquia'],
              ['Tremarctos ornatus (oso de anteojos)', 'Mamíferos', '77', 'VU + especie paraguas de los Andes'],
              ['Bos taurus BON (Blanco Orejinegro)', 'An. Domésticos', '84', 'Raza criolla antioqueña en riesgo de desplazamiento'],
          ], col_widths=[4, 2, 1.3, 5], font_size=13, header_size=14)

# SLIDE 15: Cobertura geográfica
s = content_slide('Parte 2', 'Cobertura por subregión: una brecha detectada', 15)
add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(2.2),
            [
                ('Las 158 especies propuestas están distribuidas en las 9 subregiones de Antioquia, con un mínimo de 15 especies por subregión.', 0),
                ('**Bajo Cauca y Magdalena Medio quedaron con menor cobertura**: sus ecosistemas de bosque seco y humedales están sub-representados en esta primera propuesta.', 0),
            ], size=18, space_after=18)
add_table(s, Inches(1.8), Inches(4.0), Inches(9.7), Inches(1.3),
          ['Subregión', 'Total especies propuestas', 'Comentario'],
          [
              ['Bajo Cauca', '72', 'Menor cobertura: faltan peces de ciénaga y reptiles acuáticos'],
              ['Magdalena Medio', '78', 'Menor cobertura: mismo tipo de ecosistema'],
          ], col_widths=[2.5, 2.5, 6], font_size=13, header_size=13,
          row_colors={0: ESPERA, 1: ESPERA})

# SLIDE 16: Qué le pedimos al comité
s = content_slide('Parte 2', 'Lo que le pedimos al comité hoy', 16)
add_bullets(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.3),
            [
                ('**1. Validar el puntaje de estado de conservación (P1)** contra los Libros Rojos colombianos, más restrictivos en varios casos que la IUCN global.', 0),
                ('**2. Reforzar Bajo Cauca y Magdalena Medio** con especies de ecosistemas acuáticos y de bosque seco.', 0),
                ('**3. Resolver 8 cupos adicionales** entre los grupos más ricos, o especies urgentes que el comité identifique y no estén en esta propuesta.', 0),
                ('**4. Priorizar la gestión fotográfica**: 35 de las 158 especies aún no tienen foto lista; se necesita apoyo de Instituto Humboldt, CORANTIOQUIA, ProAves y universidades.', 0),
                ('**5. Aprobar una lista de espera de 40 especies adicionales**, para reemplazos y para una Fase 2 del catálogo.', 0),
            ], size=18, space_after=18)

# SLIDE 17: Cierre
s = add_slide()
set_bg(s, DARK_GREEN)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.6), Inches(0.9), height=Inches(1.1))
add_textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1),
            'Gracias', size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT)
add_textbox(s, Inches(1), Inches(3.5), Inches(11.3), Inches(1.4),
            'Quedamos atentos a la retroalimentación del comité sobre los 5 puntos de la lámina anterior.',
            size=18, color=RGBColor(0xC8, 0xE6, 0xC9), align=PP_ALIGN.CENTER, font=FONT)
add_textbox(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
            'Antioquia Natural · Gobernación de Antioquia · Secretaría de Ambiente',
            size=13, color=RGBColor(0x9C, 0xCC, 0xB0), align=PP_ALIGN.CENTER, font=FONT)

os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT_FILE)
print(f'✓ {OUT_FILE}')
print(f'  Total de diapositivas: {len(prs.slides.__iter__.__self__._sldIdLst)}')
